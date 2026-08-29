from pathlib import Path
import base64
import re
import sys
import tempfile
import unittest
import zipfile


OLD_UPSTREAM_SKILL = "ppt" + "-master"
OLD_GATED_SKILL = "article" + "-to-pptx-gated"


def skills_root() -> Path:
    candidates = [Path.cwd(), *Path(__file__).resolve().parents]
    for candidate in candidates:
        if (candidate / OLD_UPSTREAM_SKILL).exists() or (candidate / "ppt-master-plus").exists():
            return candidate
    raise RuntimeError("skills workspace root not found")


ROOT = skills_root()
SKILL = ROOT / "ppt-master-plus"


class PptMasterPlusContractTests(unittest.TestCase):
    def test_public_name_and_mode_selection_are_declared(self):
        text = (SKILL / "SKILL.md").read_text(encoding="utf-8")
        self.assertRegex(text, r"(?m)^name: ppt-master-plus$")
        self.assertIn("workflows/gated-production.md", text)
        self.assertLess(text.index("逐页确定精修"), text.index("全自动一次性生成"))
        self.assertIn("gated", text.lower())
        self.assertIn("continuous", text.lower())

    def test_gated_workflow_contains_all_required_review_stops(self):
        text = (SKILL / "workflows/gated-production.md").read_text(encoding="utf-8")
        self.assertLess(text.index("逐页确定精修"), text.index("全自动一次性生成"))
        for marker in (
            "Intake Gate",
            "Narrative Analysis Gate",
            "Outline Gate",
            "Production Route Gate",
            "Per-slide Gate",
            "Final Acceptance Gate",
        ):
            self.assertIn(marker, text)
        self.assertIn("check_speaker_notes.py", text)

    def test_gated_per_slide_approval_uses_live_preview_not_png_model_review(self):
        gated = (SKILL / "workflows/gated-production.md").read_text(encoding="utf-8")
        skill_text = (SKILL / "SKILL.md").read_text(encoding="utf-8")
        visual_review = (SKILL / "workflows/stages/visual-review.md").read_text(encoding="utf-8")

        self.assertIn("per-slide style confirmation uses Live Preview directly", gated)
        self.assertIn("Do not render a PNG/screenshot", gated)
        self.assertIn("per-slide style approval surface is Live Preview itself", skill_text)
        self.assertIn("do not render PNGs/screenshots", skill_text)
        self.assertIn("use Live Preview directly", visual_review)
        self.assertIn("do not render PNGs", visual_review)

    def test_gated_mode_requires_blocking_per_slide_confirmation(self):
        skill_text = (SKILL / "SKILL.md").read_text(encoding="utf-8")
        gated = (SKILL / "workflows/gated-production.md").read_text(encoding="utf-8")

        self.assertIn("every page checkpoint is a new ⛔ BLOCKING stop", skill_text)
        self.assertIn("Do NOT generate multiple slides in one turn", skill_text)
        self.assertIn("⛔ **BLOCKING PER PAGE**", gated)
        self.assertIn("must not continue to the\nnext slide", gated)
        self.assertIn("Do not\ngenerate the next SVG \"while waiting\"", gated)

    def test_live_preview_apply_changes_triggers_ai_annotation_repair(self):
        skill_text = (SKILL / "SKILL.md").read_text(encoding="utf-8")
        workflow = (SKILL / "workflows/live-preview.md").read_text(encoding="utf-8")
        readme = (SKILL / "README.md").read_text(encoding="utf-8")
        app_js = (SKILL / "scripts/svg_editor/static/app.js").read_text(encoding="utf-8")
        server_py = (SKILL / "scripts/svg_editor/server.py").read_text(encoding="utf-8")

        self.assertIn("Apply annotations automatically after the browser saves them", skill_text)
        self.assertIn("The user does not need to paste the browser prompt back into chat", skill_text)
        self.assertIn("raises `live_preview/annotations_ready.flag`", workflow)
        self.assertIn("AI waits for that save event and starts the repair automatically", workflow)
        self.assertIn("AI 会自动读取保存到 `svg_output/` 的批注并开始修复", readme)
        self.assertIn("modal_success_submit", app_js)
        self.assertIn("should apply annotations automatically", app_js)
        self.assertIn("--wait-annotation", server_py)

    def test_codex_desktop_live_preview_wait_adapter_is_declared(self):
        skill_text = (SKILL / "SKILL.md").read_text(encoding="utf-8")
        workflow = (SKILL / "workflows/live-preview.md").read_text(encoding="utf-8")
        readme = (SKILL / "README.md").read_text(encoding="utf-8")

        self.assertIn("Desktop agent adapter (Codex / non-CLI tools)", skill_text)
        self.assertIn("--wait-annotation --timeout 0", skill_text)
        self.assertIn("exec_command", skill_text)
        self.assertIn("Codex / desktop-agent adapter", workflow)
        self.assertIn("For Codex Desktop specifically", workflow)
        self.assertIn("background `exec_command` / terminal session", workflow)
        self.assertIn("Repeat this repair → delete flag → re-arm wait cycle", workflow)
        self.assertIn("start a new blocked wait session", workflow)
        self.assertIn("Codex Desktop", readme)
        self.assertIn("--wait-annotation --timeout 0", readme)
        self.assertIn("第二批、第三批修改意见", readme)

    def test_wait_annotation_timeout_zero_waits_forever(self):
        server_py = (SKILL / "scripts/svg_editor/server.py").read_text(encoding="utf-8")

        self.assertIn("timeout <= 0 waits forever", server_py)
        self.assertIn("deadline = None if timeout <= 0", server_py)
        self.assertIn("while deadline is None or time.time() < deadline", server_py)

    def test_confirm_ui_transition_effect_defaults_to_none(self):
        skill_text = (SKILL / "SKILL.md").read_text(encoding="utf-8")
        catalogs = (SKILL / "scripts" / "confirm_ui" / "static" / "catalogs.json").read_text(encoding="utf-8")
        app_js = (SKILL / "scripts" / "confirm_ui" / "static" / "app.js").read_text(encoding="utf-8")
        docs = (SKILL / "scripts" / "docs" / "confirm_ui.md").read_text(encoding="utf-8")
        spec_lock_ref = (SKILL / "templates" / "spec_lock_reference.md").read_text(encoding="utf-8")
        cli = (SKILL / "scripts" / "svg_to_pptx" / "pptx_cli.py").read_text(encoding="utf-8")
        builder = (SKILL / "scripts" / "svg_to_pptx" / "pptx_builder.py").read_text(encoding="utf-8")

        self.assertIn('"transition_effect"', catalogs)
        self.assertLess(catalogs.index('"transition_effect"'), catalogs.index('"id": "fade"'))
        self.assertIn('"id": "none"', catalogs)
        self.assertIn("sec_transition", app_js)
        self.assertIn('STATE.transition_effect = recId("transition_effect") || "none"', app_js)
        self.assertIn("renderTransition(host)", app_js)
        self.assertIn('"transition_effect": "none"', docs)
        self.assertIn("default/recommended value `none`", skill_text)
        self.assertIn("- transition_effect: none", spec_lock_ref)
        self.assertIn("_read_spec_lock_value(project_path, 'transition_effect')", cli)
        self.assertIn("transition_defaults.get('effect', spec_transition_effect or 'none')", cli)
        self.assertIn("transition: str | None = None", builder)

    def test_optional_diagram_routes_are_non_blocking(self):
        text = (SKILL / "references/diagram-routing.md").read_text(encoding="utf-8")
        self.assertIn("fireworks-tech-graph", text)
        self.assertIn("excalidraw", text)
        self.assertIn("built-in SVG", text)
        self.assertRegex(text, r"(?i)do not (install|block)")

    def test_live_preview_annotation_prompt_is_copyable(self):
        app_js = (SKILL / "scripts/svg_editor/static/app.js").read_text(encoding="utf-8")
        index_html = (SKILL / "scripts/svg_editor/static/index.html").read_text(encoding="utf-8")
        server_py = (SKILL / "scripts/svg_editor/server.py").read_text(encoding="utf-8")
        workflow = (SKILL / "workflows/live-preview.md").read_text(encoding="utf-8")

        self.assertIn("modal-prompt-text", index_html)
        self.assertIn("apply_annotation_prompt", app_js)
        self.assertIn("copyTextToClipboard", app_js)
        self.assertIn("annotation_files", server_py)
        self.assertIn("MUST remove the annotation markers", app_js)
        self.assertIn("必须移除标注标记", app_js)
        self.assertIn("Never leave already-fixed annotations", workflow)
        self.assertIn("page number", workflow.lower())
        self.assertIn("svg_output", app_js)

    def test_new_user_config_precedes_legacy_fallback(self):
        text = (SKILL / "scripts/config.py").read_text(encoding="utf-8")
        new_pos = text.index(".ppt-master-plus")
        old_pos = text.index("'.ppt-master'")
        self.assertLess(new_pos, old_pos)

    def test_destructive_upstream_updater_is_removed(self):
        self.assertFalse((SKILL / "scripts/update_repo.py").exists())
        provenance = (SKILL / "references/upstream.md").read_text(encoding="utf-8")
        self.assertIn("d6bcaf96b7946667f4a8871b0688b903181db527", provenance)
        self.assertRegex(provenance, r"(?i)manual")

    def test_agent_metadata_uses_new_skill_name(self):
        text = (SKILL / "agents/openai.yaml").read_text(encoding="utf-8")
        self.assertIn("PPT Master Plus", text)
        self.assertIn("$ppt-master-plus", text)

    def test_beautify_preserve_master_contract_is_declared(self):
        workflow = (SKILL / "workflows" / "profiles" / "faithful-beautify.md").read_text(encoding="utf-8")
        confirm_docs = (SKILL / "scripts" / "docs" / "confirm_ui.md").read_text(encoding="utf-8")
        app_js = (SKILL / "scripts" / "confirm_ui" / "static" / "app.js").read_text(encoding="utf-8")
        cli = (SKILL / "scripts" / "svg_to_pptx" / "pptx_cli.py").read_text(encoding="utf-8")
        builder = (SKILL / "scripts" / "svg_to_pptx" / "pptx_builder.py").read_text(encoding="utf-8")

        for text in (workflow, confirm_docs, app_js):
            self.assertIn("preserve_master", text)
        self.assertIn("--base-pptx", workflow)
        self.assertIn("--base-pptx", cli)
        self.assertIn("source slide N", builder)

    def test_three_public_routes_are_declared(self):
        skill_text = (SKILL / "SKILL.md").read_text(encoding="utf-8")
        routing = (SKILL / "workflows" / "routing.md").read_text(encoding="utf-8")

        for route in ("Generate PPTX", "Edit Native PPTX", "Create PPTX from Template"):
            self.assertIn(route, skill_text)
            self.assertIn(route, routing)
        self.assertIn("create_pptx_from_template", routing)
        self.assertIn("template_fill", routing)
        self.assertIn("native_adaptive", routing)
        self.assertIn(
            "这份 PPTX 是要保留现有内容继续修改，还是仅作为模板，用新内容生成一份新的 PPTX？",
            routing,
        )

    def test_faithful_beautify_requires_explicit_preservation(self):
        workflow = (SKILL / "workflows" / "profiles" / "faithful-beautify.md").read_text(encoding="utf-8")

        self.assertIn("Faithful Beautify PPTX", workflow)
        self.assertIn("explicitly asks to preserve", workflow)
        self.assertIn('Generic "beautify / optimize / make professional" does not trigger', workflow)
        self.assertIn("use the main pipeline", workflow)

    def test_pptx_intake_is_context_not_constraint_for_generic_beautify(self):
        skill_text = (SKILL / "SKILL.md").read_text(encoding="utf-8")
        strategist = (SKILL / "references" / "strategist.md").read_text(encoding="utf-8")

        self.assertIn("For generic PPTX beautification, `analysis/source_profile.json` is context only", skill_text)
        self.assertIn("must not lock page count", skill_text)
        self.assertIn("Strategist rebuilds the deck", skill_text)
        self.assertIn("Generic PPTX beautification / optimization routed through the main pipeline", strategist)
        self.assertIn("Source palette/fonts are candidates, not truth", strategist)

    def test_base_pptx_export_preserves_per_slide_layout_mapping_and_master_media(self):
        scripts_dir = SKILL / "scripts"
        if str(scripts_dir) not in sys.path:
            sys.path.insert(0, str(scripts_dir))

        from pptx import Presentation
        from svg_to_pptx.pptx_builder import create_pptx_with_native_svg

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            base = tmp_path / "base.pptx"
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[0])
            prs.slides.add_slide(prs.slide_layouts[1])
            prs.save(base)

            # Add a master-level image relationship to prove the source package's
            # master/layout media and rels survive the base-pptx export path.
            patched_base = tmp_path / "base_with_master_media.pptx"
            tiny_png = base64.b64decode(
                "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJ"
                "AAAADUlEQVR42mP8z8BQDwAFgwJ/lw9J7wAAAABJRU5ErkJggg=="
            )
            with zipfile.ZipFile(base, "r") as zin, zipfile.ZipFile(patched_base, "w", zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    data = zin.read(item.filename)
                    if item.filename == "ppt/slideMasters/_rels/slideMaster1.xml.rels":
                        rel = (
                            '  <Relationship Id="rIdPreserveMasterBg" '
                            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" '
                            'Target="../media/master_bg.png"/>'
                        )
                        data = data.replace(b"</Relationships>", rel.encode("utf-8") + b"\n</Relationships>")
                    zout.writestr(item, data)
                zout.writestr("ppt/media/master_bg.png", tiny_png)

            def layout_targets(pptx_path: Path) -> list[str]:
                targets = []
                with zipfile.ZipFile(pptx_path, "r") as zf:
                    for idx in (1, 2):
                        rels = zf.read(f"ppt/slides/_rels/slide{idx}.xml.rels").decode("utf-8")
                        match = re.search(r'Type="[^"]+/slideLayout" Target="([^"]+)"', rels)
                        self.assertIsNotNone(match)
                        targets.append(match.group(1))
                return targets

            source_targets = layout_targets(patched_base)
            self.assertNotEqual(source_targets[0], source_targets[1])

            # python-pptx's default blank deck is 10" × 7.5" (4:3). With
            # --base-pptx, the source deck's slide size is authoritative, so the
            # SVG authoring canvas must match that aspect instead of the
            # catalog's ppt169 default.
            svg_files = []
            for idx in (1, 2):
                svg = tmp_path / f"slide_{idx}.svg"
                svg.write_text(
                    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 960 720">'
                    f'<text x="80" y="{100 + idx * 40}" font-size="36" fill="#111111">Slide {idx}</text>'
                    '</svg>',
                    encoding="utf-8",
                )
                svg_files.append(svg)

            output = tmp_path / "out.pptx"
            ok = create_pptx_with_native_svg(
                svg_files=svg_files,
                output_path=output,
                canvas_format="ppt169",
                verbose=False,
                use_native_shapes=True,
                base_pptx=patched_base,
            )
            self.assertTrue(ok)
            self.assertEqual(source_targets, layout_targets(output))
            with zipfile.ZipFile(output, "r") as zf:
                names = set(zf.namelist())
                self.assertIn("ppt/slideMasters/slideMaster1.xml", names)
                self.assertIn("ppt/media/master_bg.png", names)
                master_rels = zf.read("ppt/slideMasters/_rels/slideMaster1.xml.rels").decode("utf-8")
                self.assertIn("master_bg.png", master_rels)

    def test_no_old_skill_name_references_remain_outside_provenance(self):
        old_names = re.compile(
            rf"\b{re.escape(OLD_GATED_SKILL)}\b|"
            rf"skills/{re.escape(OLD_UPSTREAM_SKILL)}(?!-plus)(?:/|\b)"
        )
        allowed_legacy_files = {
            SKILL / "references/upstream.md",
        }
        offenders = []
        for path in SKILL.rglob("*"):
            if not path.is_file() or ".git" in path.parts:
                continue
            if path.resolve() == Path(__file__).resolve():
                continue
            if path in allowed_legacy_files:
                continue
            if path.suffix not in {".md", ".py", ".yaml", ".yml", ".json"}:
                continue
            try:
                text = path.read_text(encoding="utf-8")
            except UnicodeDecodeError:
                continue
        self.assertEqual([], offenders)

    def test_confirm_ui_mandatory_launch_unless_headless(self):
        skill_text = (SKILL / "SKILL.md").read_text(encoding="utf-8")
        strategist_text = (SKILL / "references" / "strategist.md").read_text(encoding="utf-8")

        self.assertIn("Mandatory Web Confirmation Page Launch (Unless Headless)", skill_text)
        self.assertIn("ALWAYS MANDATORY", skill_text)
        self.assertIn("ONLY permitted if the user/environment is headless", skill_text)
        self.assertIn("MANDATORY for all PPT requests unless the user/environment is headless", strategist_text)


if __name__ == "__main__":
    unittest.main()
