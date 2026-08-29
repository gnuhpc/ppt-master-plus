import hashlib
import json
from pathlib import Path
import shutil
import subprocess
import sys
import tempfile
import unittest
import zipfile


SKILL = Path(__file__).resolve().parents[2]
SCRIPTS = SKILL / "scripts"
VENDOR = SKILL / "vendor" / "ppt-master"


class UnifiedCapabilityContractTests(unittest.TestCase):
    def test_frozen_plus_ui_and_preview_surface_is_byte_identical(self):
        manifest = Path(__file__).with_name("frozen_surface.sha256")
        for line in manifest.read_text(encoding="utf-8").splitlines():
            expected, relative = line.split(maxsplit=1)
            actual = hashlib.sha256((SKILL / relative).read_bytes()).hexdigest()
            self.assertEqual(expected, actual, relative)

    def test_vendored_core_has_static_provenance_and_no_forbidden_surface(self):
        self.assertTrue((VENDOR / "LICENSE").is_file())
        self.assertTrue((VENDOR / "MANIFEST.sha256").is_file())
        for relative in (
            "scripts/confirm_ui",
            "scripts/svg_editor",
            "scripts/update_repo.py",
            "scripts/attribution_guard.py",
        ):
            self.assertFalse((VENDOR / relative).exists(), relative)

    def test_vendored_manifest_matches_files(self):
        for line in (VENDOR / "MANIFEST.sha256").read_text(encoding="utf-8").splitlines():
            expected, relative = line.split(maxsplit=1)
            path = VENDOR / relative
            self.assertTrue(path.is_file(), relative)
            self.assertEqual(expected, hashlib.sha256(path.read_bytes()).hexdigest(), relative)

    def test_public_entrypoints_use_one_engine_without_a_dispatcher(self):
        self.assertFalse((SCRIPTS / "ppt_master_core.py").exists())
        self.assertFalse(list(SCRIPTS.glob("*_legacy.py")))
        public_entrypoints = (
            "svg_to_pptx.py",
            "project_manager.py",
            "svg_quality_checker.py",
            "pptx_to_svg.py",
            "pptx_intake.py",
            "finalize_svg.py",
            "svg_authoring_view.py",
        )
        for entrypoint in public_entrypoints:
            text = (SCRIPTS / entrypoint).read_text(encoding="utf-8")
            self.assertNotIn("ppt_master_core", text, entrypoint)
            self.assertIn("vendor", text, entrypoint)
        self.assertTrue((SCRIPTS / "native_preview.py").is_file())
        for removed in (
            "native_enhance_pptx.py",
            "native_narration_pptx.py",
            "template_fill_pptx.py",
            "template_fill_pptx",
        ):
            self.assertFalse((SCRIPTS / removed).exists(), removed)

        export_help = subprocess.run(
            [sys.executable, str(SCRIPTS / "svg_to_pptx.py"), "--help"],
            capture_output=True,
            text=True,
        )
        self.assertEqual(0, export_help.returncode, export_help.stderr)
        self.assertIn("--native-charts-and-tables", export_help.stdout)
        self.assertIn("--pptx-structure", export_help.stdout)
        self.assertNotIn("--master-core", export_help.stdout)

        project_help = subprocess.run(
            [sys.executable, str(SCRIPTS / "project_manager.py"), "--help"],
            capture_output=True,
            text=True,
        )
        self.assertEqual(0, project_help.returncode, project_help.stderr)
        self.assertIn("page-context", project_help.stdout)
        self.assertIn("import-template", project_help.stdout)
        self.assertIn("set-output", project_help.stdout)
        self.assertIn("configure-native", project_help.stdout)
        self.assertIn("native-plan", project_help.stdout)

    def test_quick_generate_is_rejected_at_public_boundary(self):
        result = subprocess.run(
            [sys.executable, str(SCRIPTS / "svg_to_pptx.py"), "demo", "--quick-generate"],
            capture_output=True,
            text=True,
        )
        self.assertNotEqual(0, result.returncode)
        self.assertIn("does not expose Quick Generate", result.stderr + result.stdout)

    def test_template_indexes_are_valid_and_payload_keys_unique(self):
        cases = {
            "brands": None,
            "styles": None,
            "layouts": None,
            "decks": None,
            "charts": "charts",
            "tables": "tables",
            "sounds": "sounds",
        }
        for family, payload_key in cases.items():
            index = SKILL / "templates" / family / f"{family}_index.json"
            duplicate_keys = []

            def checked_object(pairs):
                result = {}
                for key, value in pairs:
                    if key in result:
                        duplicate_keys.append(key)
                    result[key] = value
                return result

            payload = json.loads(
                index.read_text(encoding="utf-8"), object_pairs_hook=checked_object
            )
            self.assertEqual([], duplicate_keys, family)
            entries = payload[payload_key] if payload_key else payload
            self.assertTrue(entries, family)
            if isinstance(entries, list):
                ids = [entry["id"] for entry in entries]
                self.assertEqual(len(ids), len(set(ids)), family)
                for entry in entries:
                    self.assertTrue((index.parent / entry["file"]).is_file(), entry["file"])
            else:
                aliases = payload.get("aliases", {}) if payload_key else {}
                self.assertFalse(set(entries) & set(aliases), family)
                for entry_id in entries:
                    self.assertTrue(
                        (index.parent / entry_id).exists()
                        or (index.parent / f"{entry_id}.svg").is_file(),
                        f"{family}:{entry_id}",
                    )

    def test_unified_helpers_are_public_and_quick_help_does_not_advertise_router(self):
        helpers = (
            "source_to_md.py",
            "preset_shape_svg.py",
            "shape_boolean_svg.py",
            "pptx_delivery_check.py",
            "pptx_opc_validation.py",
            "sound_sync.py",
            "video_motion_plan.py",
            "video_sound_mix.py",
            "video_subtitles.py",
            "visualization_recall.py",
        )
        for helper in helpers:
            self.assertTrue((SCRIPTS / helper).is_file(), helper)
        result = subprocess.run(
            [sys.executable, str(SCRIPTS / "svg_to_pptx.py"), "--help"],
            capture_output=True,
            text=True,
        )
        self.assertEqual(0, result.returncode, result.stderr)
        self.assertNotIn("--quick-generate", result.stdout)

        visualization = subprocess.run(
            [
                sys.executable,
                str(SCRIPTS / "visualization_recall.py"),
                "validate",
                "chart/area_chart",
                "table/comparison_matrix",
            ],
            capture_output=True,
            text=True,
        )
        self.assertEqual(0, visualization.returncode, visualization.stderr + visualization.stdout)
        validated = json.loads(visualization.stdout)
        self.assertEqual([], validated["invalid"])
        self.assertEqual(2, len(validated["valid"]))

    def test_preset_registry_exposes_all_187_native_shapes(self):
        result = subprocess.run(
            [sys.executable, str(SCRIPTS / "preset_shape_svg.py"), "list"],
            capture_output=True,
            text=True,
        )
        self.assertEqual(0, result.returncode, result.stderr)
        presets = [line for line in result.stdout.splitlines() if line.strip()]
        self.assertEqual(187, len(presets))
        self.assertIn("rect", presets)
        self.assertIn("actionButtonHome", presets)

        boolean_fixture = Path(__file__).with_name("fixtures") / "boolean_shapes.svg"
        merged = subprocess.run(
            [
                sys.executable,
                str(SCRIPTS / "shape_boolean_svg.py"),
                "render",
                str(boolean_fixture),
                "--operation",
                "union",
                "--source",
                "left",
                "--source",
                "right",
                "--id",
                "merged-shape",
            ],
            capture_output=True,
            text=True,
        )
        self.assertEqual(0, merged.returncode, merged.stderr)
        self.assertIn('id="merged-shape"', merged.stdout)
        self.assertIn("<path", merged.stdout)

    def test_native_chart_table_formula_and_hyperlink_fixture_passes_core_checker(self):
        fixture = Path(__file__).with_name("fixtures") / "native_objects.svg"
        result = subprocess.run(
            [
                sys.executable,
                str(SCRIPTS / "svg_quality_checker.py"),
                str(fixture),
                "--format",
                "ppt169",
            ],
            capture_output=True,
            text=True,
        )
        self.assertEqual(0, result.returncode, result.stderr + result.stdout)
        for receipt in ("charts 1", "tables 1", "formulas 1"):
            self.assertIn(receipt, result.stdout)
        source = fixture.read_text(encoding="utf-8")
        self.assertIn('href="https://example.com/report"', source)

    def test_native_fixture_exports_real_drawingml_objects(self):
        fixture = Path(__file__).with_name("fixtures") / "native_objects.svg"
        with tempfile.TemporaryDirectory() as tmp:
            project = Path(tmp) / "project"
            (project / "svg_output").mkdir(parents=True)
            (project / "exports").mkdir()
            shutil.copy2(fixture, project / "svg_output" / "01_native.svg")
            (project / "spec_lock.md").write_text(
                """<!-- ppt-master-schema: spec-lock/v1 -->
# Execution Lock
## canvas
- viewBox: 0 0 1280 720
- format: PPT 16:9
## typography
- font_family: Arial
- title_family: Arial
- body_family: Arial
- body: 18
- title: 34
## colors
- bg: #FFFFFF
- primary: #173B57
- accent: #04CF82
- text: #173B57
## page_rhythm
- P01: native objects fixture
## pptx_structure
- mode: flat
""",
                encoding="utf-8",
            )
            checked = subprocess.run(
                [
                    sys.executable,
                    str(SCRIPTS / "svg_quality_checker.py"),
                    str(project),
                    "--stage",
                    "final",
                    "--json",
                ],
                capture_output=True,
                text=True,
            )
            self.assertEqual(0, checked.returncode, checked.stderr + checked.stdout)
            output = project / "exports" / "native-fixture.pptx"
            exported = subprocess.run(
                [
                    sys.executable,
                    str(SCRIPTS / "svg_to_pptx.py"),
                    str(project),
                    "--native-charts-and-tables",
                    "--conversion-trace",
                    "-t",
                    "fade",
                    "-a",
                    "entrance_fade",
                    "-o",
                    str(output),
                ],
                capture_output=True,
                text=True,
            )
            self.assertEqual(0, exported.returncode, exported.stderr + exported.stdout)
            with zipfile.ZipFile(output) as package:
                names = package.namelist()
                slide = package.read("ppt/slides/slide1.xml")
                self.assertTrue(any(name.startswith("ppt/charts/chart") for name in names))
                self.assertTrue(any(name.startswith("ppt/embeddings/") for name in names))
                self.assertIn(b"<a:tbl>", slide)
                self.assertIn(b"<m:oMath", slide)
                self.assertIn(b"hlinkClick", slide)
                self.assertIn(b"<p:transition", slide)
                self.assertIn(b"<p:timing", slide)

            trace = project / "validation" / "native-fixture.trace.json"
            motion_plan = project / "validation" / "native-fixture.video-motion.json"
            planned = subprocess.run(
                [
                    sys.executable,
                    str(SCRIPTS / "video_motion_plan.py"),
                    str(trace),
                    "-o",
                    str(motion_plan),
                ],
                capture_output=True,
                text=True,
            )
            self.assertEqual(0, planned.returncode, planned.stderr + planned.stdout)
            plan = json.loads(motion_plan.read_text(encoding="utf-8"))
            self.assertEqual(1, plan["slide_count"])
            self.assertGreater(plan["object_count"], 0)

    def test_layered_roundtrip_and_delivery_postflight_fixture(self):
        source = SKILL / "harness" / "testing" / "original-sample.pptx"
        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "roundtrip"
            converted = subprocess.run(
                [
                    sys.executable,
                    str(SCRIPTS / "pptx_to_svg.py"),
                    str(source),
                    "--inheritance-mode",
                    "both",
                    "--roundtrip",
                    "-o",
                    str(output),
                ],
                capture_output=True,
                text=True,
            )
            self.assertEqual(0, converted.returncode, converted.stderr)
            report = json.loads((output / "validation" / "conversion-report.json").read_text())
            self.assertEqual(6, report["summary"]["slides"])
            self.assertTrue((output / "analysis" / "native_structure.json").is_file())
            self.assertTrue((output / "sources" / "source.pptx").is_file())
            authoring = "\n".join(
                path.read_text(encoding="utf-8")
                for path in sorted((output / "authoring-svg-flat").glob("slide_*.svg"))
            )
            for marker in (
                "data-pptx-source-ref",
                "data-pptx-object",
                "data-pptx-prst",
                "data-pptx-text-model",
            ):
                self.assertIn(marker, authoring)

            identity = output / "exports" / "identity.pptx"
            exported = subprocess.run(
                [
                    sys.executable,
                    str(SCRIPTS / "svg_to_pptx.py"),
                    str(output),
                    "--roundtrip",
                    "-o",
                    str(identity),
                ],
                capture_output=True,
                text=True,
            )
            self.assertEqual(0, exported.returncode, exported.stderr + exported.stdout)
            with zipfile.ZipFile(source) as original, zipfile.ZipFile(identity) as restored:
                for index in range(1, 7):
                    for relative in (
                        f"ppt/slides/slide{index}.xml",
                        f"ppt/slides/_rels/slide{index}.xml.rels",
                    ):
                        self.assertEqual(original.read(relative), restored.read(relative), relative)

        opc = subprocess.run(
            [sys.executable, str(SCRIPTS / "pptx_opc_validation.py"), str(source), "--json"],
            capture_output=True,
            text=True,
        )
        self.assertEqual(0, opc.returncode, opc.stderr)
        self.assertTrue(json.loads(opc.stdout)["ok"])

        delivery = subprocess.run(
            [sys.executable, str(SCRIPTS / "pptx_delivery_check.py"), str(source)],
            capture_output=True,
            text=True,
        )
        self.assertEqual(0, delivery.returncode, delivery.stderr)
        self.assertIn(json.loads(delivery.stdout)["status"], {"passed", "passed-with-advisories"})

    def test_public_workflow_surface_has_only_three_routes(self):
        workflows = SKILL / "workflows"
        for route in ("generate-pptx.md", "edit-native-pptx.md", "create-pptx-from-template.md"):
            self.assertTrue((workflows / route).is_file(), route)
        for removed in (
            "native-enhance-pptx.md",
            "native-narration-pptx.md",
            "template-fill-pptx.md",
            "create-brand.md",
            "create-template.md",
            "beautify-pptx.md",
        ):
            self.assertFalse((workflows / removed).exists(), removed)
        self.assertTrue((workflows / "profiles" / "faithful-beautify.md").is_file())
        self.assertTrue((workflows / "internal" / "template-authoring" / "create-brand.md").is_file())
        routing = (workflows / "routing.md").read_text(encoding="utf-8")
        template_route = (workflows / "create-pptx-from-template.md").read_text(
            encoding="utf-8"
        )
        self.assertIn("visualization_style_source: skill_builtin", routing)
        self.assertIn("template_origin: user_provided_pptx", routing)
        self.assertIn("templates/charts/charts_index.json", template_route)
        self.assertIn("templates/tables/tables_index.json", template_route)
        self.assertIn("template-native chart or", template_route)
        self.assertIn("table formatting is never style authority", template_route)
        self.assertIn("explicitly provided by the", template_route)
        self.assertIn("Reject `.potx`", template_route)

    def test_template_intake_infers_ordinary_slots_and_publishes_four_contracts(self):
        from pptx import Presentation
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "synthetic-template.pptx"
            prs = Presentation()
            first = prs.slides.add_slide(prs.slide_layouts[6])
            first.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text = (
                "字体规范 Microsoft YaHei #112233"
            )
            second = prs.slides.add_slide(prs.slide_layouts[6])
            second.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text = (
                "这里是标题，限制一行"
            )
            prs.save(source)

            analysis = root / "analysis"
            result = subprocess.run(
                [
                    sys.executable,
                    str(SCRIPTS / "pptx_intake.py"),
                    str(source),
                    "-o",
                    str(analysis),
                    "--intent",
                    "template",
                    "--no-previews",
                ],
                capture_output=True,
                text=True,
            )
            self.assertEqual(0, result.returncode, result.stderr + result.stdout)
            names = (
                "template_manifest.json",
                "template_design_tokens.json",
                "template_archetypes.json",
                "template_assets.json",
            )
            for name in names:
                self.assertTrue((analysis / name).is_file(), name)
            manifest = json.loads((analysis / names[0]).read_text(encoding="utf-8"))
            tokens = json.loads((analysis / names[1]).read_text(encoding="utf-8"))
            archetypes = json.loads((analysis / names[2]).read_text(encoding="utf-8"))
            self.assertEqual(1, manifest["page_role_counts"]["guide"])
            policy = manifest["visualization_style_policy"]
            self.assertEqual("skill_builtin", policy["source"])
            self.assertFalse(policy["template_visualization_style_reusable"])
            self.assertEqual("skill_builtin", tokens["visualization_style_source"])
            self.assertIn("chart_style", tokens["excluded_template_style_domains"])
            self.assertEqual(
                "skill_builtin", archetypes["visualization_policy"]["style_source"]
            )
            self.assertEqual(
                "content_sources",
                archetypes["visualization_policy"]["content_and_data_source"],
            )
            ordinary = archetypes["slides"][1]["objects"][0]
            self.assertIsNone(ordinary["placeholder"])
            self.assertTrue(ordinary["editable"])
            self.assertEqual(1, ordinary["capacity"]["max_lines"])

    def test_project_template_import_is_read_only_and_records_separate_roles(self):
        source = SKILL / "harness" / "testing" / "original-sample.pptx"
        source_hash = hashlib.sha256(source.read_bytes()).hexdigest()
        with tempfile.TemporaryDirectory() as tmp:
            project = Path(tmp) / "project"
            (project / "sources").mkdir(parents=True)
            (project / "sources" / "content.md").write_text("# New content\n", encoding="utf-8")
            result = subprocess.run(
                [
                    sys.executable,
                    str(SCRIPTS / "project_manager.py"),
                    "import-template",
                    str(project),
                    str(source),
                    "--no-previews",
                ],
                capture_output=True,
                text=True,
            )
            self.assertEqual(0, result.returncode, result.stderr + result.stdout)
            copied = project / "template" / "source.pptx"
            self.assertTrue(copied.is_file())
            self.assertEqual(source_hash, hashlib.sha256(source.read_bytes()).hexdigest())
            self.assertEqual(source_hash, hashlib.sha256(copied.read_bytes()).hexdigest())
            context = json.loads((project / "route_context.json").read_text(encoding="utf-8"))
            self.assertEqual("create_pptx_from_template", context["route"])
            self.assertEqual("template_fill", context["profile"])
            self.assertEqual("skill_builtin", context["visualization_style_source"])
            self.assertIn(
                "templates/charts/charts_index.json", context["visualization_catalogs"]
            )
            self.assertEqual("user_provided_pptx", context["template_origin"])
            self.assertTrue(context["template_import_required"])
            self.assertEqual("readonly_copy", context["template_import_mode"])
            self.assertEqual("template/source.pptx", context["template_pptx"])
            self.assertEqual(["sources/content.md"], context["content_sources"])
            self.assertIsNone(context["output_pptx"])

    def test_template_import_rejects_potx_and_requires_user_pptx_file(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            project = root / "project"
            project.mkdir()
            potx = root / "template.potx"
            potx.write_bytes(b"not-a-pptx")
            result = subprocess.run(
                [
                    sys.executable,
                    str(SCRIPTS / "project_manager.py"),
                    "import-template",
                    str(project),
                    str(potx),
                    "--no-previews",
                ],
                capture_output=True,
                text=True,
            )
            self.assertNotEqual(0, result.returncode)
            self.assertIn("Template must be a .pptx package", result.stderr + result.stdout)

    def test_native_preview_adapter_never_replaces_nonempty_generate_workspace(self):
        with tempfile.TemporaryDirectory() as tmp:
            project = Path(tmp) / "project"
            authoring = project / "authoring-svg-flat"
            authoring.mkdir(parents=True)
            (authoring / "slide_01.svg").write_text("<svg/>", encoding="utf-8")
            prepared = subprocess.run(
                [sys.executable, str(SCRIPTS / "native_preview.py"), "prepare", str(project)],
                capture_output=True,
                text=True,
            )
            self.assertEqual(0, prepared.returncode, prepared.stderr)
            self.assertTrue((project / "svg_output").is_symlink())
            cleaned = subprocess.run(
                [sys.executable, str(SCRIPTS / "native_preview.py"), "clean", str(project)],
                capture_output=True,
                text=True,
            )
            self.assertEqual(0, cleaned.returncode, cleaned.stderr)
            self.assertFalse((project / "svg_output").exists())

            (project / "svg_output").mkdir()
            (project / "svg_output" / "owned.svg").write_text("<svg/>", encoding="utf-8")
            refused = subprocess.run(
                [sys.executable, str(SCRIPTS / "native_preview.py"), "prepare", str(project)],
                capture_output=True,
                text=True,
            )
            self.assertNotEqual(0, refused.returncode)
            self.assertTrue((project / "svg_output" / "owned.svg").is_file())

    def test_unified_native_modules_and_page_plan_clone_are_fail_closed(self):
        with tempfile.TemporaryDirectory() as tmp:
            project = Path(tmp) / "project"
            authoring = project / "authoring-svg-flat"
            authoring.mkdir(parents=True)
            (authoring / "slide_01.svg").write_text("<svg/>", encoding="utf-8")
            configured = subprocess.run(
                [
                    sys.executable,
                    str(SCRIPTS / "project_manager.py"),
                    "configure-native",
                    str(project),
                    "--module",
                    "page_plan",
                    "--module",
                    "content_edit",
                ],
                capture_output=True,
                text=True,
            )
            self.assertEqual(0, configured.returncode, configured.stderr + configured.stdout)
            context = json.loads((project / "route_context.json").read_text(encoding="utf-8"))
            self.assertEqual("edit_native_pptx", context["route"])
            self.assertEqual(["page_plan", "content_edit"], context["modules"])

            planned = subprocess.run(
                [
                    sys.executable,
                    str(SCRIPTS / "project_manager.py"),
                    "native-plan",
                    str(project),
                    "--page",
                    "1:first.svg",
                    "--page",
                    "1:second.svg",
                ],
                capture_output=True,
                text=True,
            )
            self.assertEqual(0, planned.returncode, planned.stderr + planned.stdout)
            plan = json.loads((project / "page_plan.json").read_text(encoding="utf-8"))
            self.assertEqual(["first.svg", "second.svg"], [row["svg"] for row in plan["pages"]])
            self.assertTrue((authoring / "first.svg").is_file())
            self.assertTrue((authoring / "second.svg").is_file())

            refused = subprocess.run(
                [
                    sys.executable,
                    str(SCRIPTS / "project_manager.py"),
                    "configure-native",
                    str(project),
                    "--module",
                    "visible_content_locked",
                    "--module",
                    "content_edit",
                ],
                capture_output=True,
                text=True,
            )
            self.assertNotEqual(0, refused.returncode)


if __name__ == "__main__":
    unittest.main()
